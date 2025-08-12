# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

import pandas as pd
import re
from datetime import date
from calendar import monthrange

# -----------------------
# SLIDE & LAYOUT
# -----------------------
SLIDE_W_IN = 20.0
SLIDE_H_IN = 9.0

TOP_PAD_IN    = 0.40
TITLE_H_IN    = 0.40
LEGEND_H_IN   = 0.60
HEADER_H_IN   = 1.00
BOTTOM_PAD_IN = 0.09
LEFT_PAD_IN   = 0.40
RIGHT_PAD_IN  = 0.40

TYPE_COL_W_IN = 1.50
WORK_COL_W_IN = 2.50
SIDEBAR_W_IN  = TYPE_COL_W_IN + WORK_COL_W_IN

MONTHS = 12
ROWS_PER_SLIDE = 20  # exact pagination as requested

# Colors
BLUE_HDR    = RGBColor(91,155,213)
MONTH_ODD   = RGBColor(224,242,255)
MONTH_EVEN  = RGBColor(190,220,240)
WHITE       = RGBColor(255,255,255)
BLACK       = RGBColor(0,0,0)
NAVY        = RGBColor(0,0,128)
TODAY_GRN   = RGBColor(0,176,80)

STATUS_COLORS = {
    "On Track":  RGBColor(0,176,80),
    "At Risk":   RGBColor(255,192,0),
    "Off Track": RGBColor(255,0,0),
    "Complete":  RGBColor(0,112,192),
    "TBC":       RGBColor(191,191,191),
}

# Milestone shapes
CIRCLE_SIZE_IN = 0.30
STAR_SIZE_IN   = 0.40  # Major milestone

# Smart label tuning
LABEL_W_IN        = 2.6
LABEL_LINE_H_IN   = 0.18
LABEL_X_GAP_IN    = 0.10
LABEL_Y_OFFSET_IN = 0.12
LANES_PER_SIDE    = 4
CLUSTER_X_FRACT   = 0.25
LAST4_START_IDX   = 8   # Sep..Dec

# Short-hands
def IN(v): return Inches(v)

# -----------------------
# TEXT HELPERS (sorting & labels)
# -----------------------
def clean_text(s: str) -> str:
    if pd.isna(s): return ""
    return str(s).strip()

def type_bucket(s: str) -> str:
    s = clean_text(s).casefold()
    m = re.match(r"([a-z0-9]+)", re.sub(r"[^a-z0-9]", "", s))
    return m.group(1) if m else s

def wrap_by_words(text: str, n_words: int=3) -> str:
    words = clean_text(text).split()
    if not words: return ""
    return "\n".join(" ".join(words[i:i+n_words]) for i in range(0, len(words), n_words))

def est_label_h(text: str) -> float:
    lines = max(1, text.count("\n")+1)
    return lines * LABEL_LINE_H_IN

def rects_intersect(a, b):
    ax0, ay0, ax1, ay1 = a
    bx0, by0, bx1, by1 = b
    return (ax0 < bx1) and (bx0 < ax1) and (ay0 < by1) and (by0 < ay1)

def clamp(v, lo, hi):
    return max(lo, min(hi, v))

# -----------------------
# TABLE BORDERS (OXML)  — Option A
# -----------------------
def set_cell_border(cell, color_hex="FFFFFF", width="6350"):
    """
    True cell borders via OXML (python-pptx has no high-level API).
    color_hex without '#'; width ~ 6350 thin, 12700 thicker.
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("L", "R", "T", "B"):
        ln = OxmlElement(f'a:ln{edge}')
        ln.set('w', width)
        solid = OxmlElement('a:solidFill')
        rgb = OxmlElement('a:srgbClr'); rgb.set('val', color_hex)
        solid.append(rgb); ln.append(solid)
        dash = OxmlElement('a:prstDash'); dash.set('val','solid')
        ln.append(dash)
        tcPr.append(ln)

# -----------------------
# LEGEND & TITLE
# -----------------------
def add_title(slide, title_text):
    tb = slide.shapes.add_textbox(IN(LEFT_PAD_IN), IN(TOP_PAD_IN), IN(SLIDE_W_IN-LEFT_PAD_IN-RIGHT_PAD_IN), IN(TITLE_H_IN))
    p = tb.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24); p.font.bold = True

def add_full_width_legend(slide):
    items = [
        ("Major Milestone", MSO_SHAPE.STAR_5_POINT, STATUS_COLORS["On Track"]),
        ("On Track",  MSO_SHAPE.OVAL, STATUS_COLORS["On Track"]),
        ("At Risk",   MSO_SHAPE.OVAL, STATUS_COLORS["At Risk"]),
        ("Off Track", MSO_SHAPE.OVAL, STATUS_COLORS["Off Track"]),
        ("Complete",  MSO_SHAPE.OVAL, STATUS_COLORS["Complete"]),
        ("TBC",       MSO_SHAPE.OVAL, STATUS_COLORS["TBC"]),
    ]
    slot = (SLIDE_W_IN - LEFT_PAD_IN - RIGHT_PAD_IN)/len(items)
    y = TOP_PAD_IN + TITLE_H_IN/4.0
    for i,(label, shp_kind, col) in enumerate(items):
        x = LEFT_PAD_IN + i*slot + 0.2
        s = slide.shapes.add_shape(shp_kind, IN(x), IN(y), IN(0.30), IN(0.30))
        s.fill.solid(); s.fill.fore_color.rgb = col
        s.line.fill.background()
        tb = slide.shapes.add_textbox(IN(x+0.40), IN(y-0.02), IN(2.5), IN(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = label; p.font.size = Pt(14)

# -----------------------
# BUILD EDITABLE TABLE  (and return geometry for shapes)
# -----------------------
def build_table(slide, year, groups):
    """
    Returns geometry dict:
      {
        'left': left_origin_in, 'top': top_origin_in,
        'first_row_top': first_row_top_in,
        'dates_left': dates_left_in,
        'month_w': month_w_in,
        'row_h': row_h_in,
        'row_count': int
      }
    """
    # overall canvas for table area
    left_in = LEFT_PAD_IN
    top_in  = TOP_PAD_IN + TITLE_H_IN + LEGEND_H_IN
    total_w_in = SLIDE_W_IN - LEFT_PAD_IN - RIGHT_PAD_IN
    total_h_in = SLIDE_H_IN - (TOP_PAD_IN + TITLE_H_IN + LEGEND_H_IN) - BOTTOM_PAD_IN

    row_count = max(1, len(groups))

    # column widths
    dates_w_in = total_w_in - SIDEBAR_W_IN
    month_w_in = dates_w_in / MONTHS

    # row heights: stretch to fill, always down to 0.09" from bottom
    header_h_in = HEADER_H_IN
    body_h_in   = max(0.1, total_h_in - header_h_in)
    row_h_in    = body_h_in / row_count

    # make the actual table
    rows = row_count + 1
    cols = 2 + MONTHS
    tbl_shape = slide.shapes.add_table(int(rows), int(cols), IN(left_in), IN(top_in), IN(total_w_in), IN(total_h_in))
    tbl = tbl_shape.table

    # widths
    tbl.columns[0].width = IN(TYPE_COL_W_IN)
    tbl.columns[1].width = IN(WORK_COL_W_IN)
    for c in range(MONTHS):
        tbl.columns[2+c].width = IN(month_w_in)

    # heights
    tbl.rows[0].height = IN(header_h_in)
    for r in range(1, rows):
        tbl.rows[r].height = IN(row_h_in)

    # headers: Type & Workstream (white text + borders)
    for c,title in [(0,"Type"), (1,"Workstream")]:
        cell = tbl.cell(0,c)
        cell.fill.solid(); cell.fill.fore_color.rgb = BLUE_HDR
        tf = cell.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = title; p.font.bold = True; p.font.size = Pt(18); p.font.color.rgb = WHITE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p.alignment = PP_ALIGN.CENTER
        set_cell_border(cell, "FFFFFF", "6350")

    # month headers
    for m in range(MONTHS):
        cell = tbl.cell(0, 2+m)
        cell.fill.solid(); cell.fill.fore_color.rgb = BLUE_HDR
        tf = cell.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = date(year, m+1, 1).strftime("%b %y")
        p.font.bold = True; p.font.size = Pt(18); p.font.color.rgb = WHITE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p.alignment = PP_ALIGN.CENTER
        set_cell_border(cell, "FFFFFF", "6350")

    # body rows
    for r, grp in enumerate(groups, start=1):
        t, w = grp.split("\n",1)
        # left two cells (black text) + alternating fill
        for c, txt in [(0, t), (1, w.strip("()"))]:
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = MONTH_ODD if (r % 2 == 1) else MONTH_EVEN
            tf = cell.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(15); p.font.color.rgb = BLACK
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p.alignment = PP_ALIGN.CENTER
            set_cell_border(cell, "FFFFFF", "6350")
        # date area cells + borders
        for c in range(MONTHS):
            cell = tbl.cell(r, 2+c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = MONTH_ODD if (r % 2 == 1) else MONTH_EVEN
            set_cell_border(cell, "FFFFFF", "6350")

    # geometry for placing connectors & shapes
    first_row_top_in = top_in + header_h_in
    dates_left_in    = left_in + SIDEBAR_W_IN

    # Navy center lines (thin, across dates area only)
    x_left  = dates_left_in
    x_right = dates_left_in + MONTHS*month_w_in
    for r in range(row_count):
        y_c = first_row_top_in + r*row_h_in + row_h_in/2.0
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x_left), IN(y_c), IN(x_right), IN(y_c))
        ln.line.fill.solid(); ln.line.fore_color.rgb = NAVY; ln.line.width = Pt(0.5)

    return {
        'left': left_in,
        'top': top_in,
        'first_row_top': first_row_top_in,
        'dates_left': dates_left_in,
        'month_w': month_w_in,
        'row_h': row_h_in,
        'row_count': row_count
    }

# -----------------------
# TODAY LINE
# -----------------------
def add_today_line(slide, year, geom):
    today = date.today()
    if today.year != year:
        return
    days = monthrange(today.year, today.month)[1]
    frac = (today.day - 1)/(days - 1) if days > 1 else 0.0
    x = geom['dates_left'] + (today.month - 1 + frac)*geom['month_w']
    top = geom['first_row_top']
    bottom = geom['first_row_top'] + geom['row_count']*geom['row_h']
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x), IN(top), IN(x), IN(bottom))
    conn.line.fill.solid(); conn.line.fore_color.rgb = TODAY_GRN; conn.line.width = Pt(2)
    conn.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT

# -----------------------
# MILESTONES (with smart labels)
# -----------------------
def place_labels_for_row(slide, items, geom, month_idx):
    """
    items: [{'x','y_center','title'}] for one row
    Collision-aware placement near marker; clamp inside row and dates area.
    """
    x_min = geom['dates_left']
    x_max = geom['dates_left'] + MONTHS*geom['month_w']
    row_top = items[0]['y_center'] - geom['row_h']/2.0
    row_h   = geom['row_h']
    col_w   = geom['month_w']

    placed_rects = []
    last_x = None
    alt_above = True

    for it in items:
        x = it['x']; y_center = it['y_center']
        raw = it['title']

        prefer_above = False
        # for last 4 months -> above + wrap
        if month_idx >= LAST4_START_IDX:
            prefer_above = True
            label_text = wrap_by_words(raw, 3)
        else:
            label_text = raw
            # right overflow?
            if x + LABEL_X_GAP_IN + LABEL_W_IN > x_max:
                # can we place left?
                if x - LABEL_X_GAP_IN - LABEL_W_IN >= x_min:
                    prefer_above = False
                else:
                    prefer_above = True
                    label_text = wrap_by_words(raw, 3)

        label_h = est_label_h(label_text)

        # cluster detection (if close in x within same row)
        if last_x is not None and abs(x - last_x) < CLUSTER_X_FRACT*col_w:
            alt_above = not alt_above
        last_x = x

        # base position candidate
        if prefer_above:
            bx = clamp(x - LABEL_W_IN/2.0, x_min, x_max - LABEL_W_IN)
            by = y_center - LABEL_Y_OFFSET_IN - label_h
            lane_order = [0] + [i for k in range(1, LANES_PER_SIDE+1) for i in (k, -k)]
        else:
            place_right = (x + LABEL_X_GAP_IN + LABEL_W_IN) <= x_max
            base_y = y_center - (LABEL_Y_OFFSET_IN if alt_above else -LABEL_Y_OFFSET_IN)
            bx = clamp(x + (LABEL_X_GAP_IN if place_right else -LABEL_X_GAP_IN - LABEL_W_IN), x_min, x_max - LABEL_W_IN)
            by = base_y
            lane_order = [0] + [i for k in range(1, LANES_PER_SIDE+1) for i in (-k, k)]

        lane_step = min(LABEL_LINE_H_IN * 1.05, max(0.05, (row_h/2.0) - LABEL_LINE_H_IN))
        y_lo = row_top + 0.04
        y_hi = row_top + row_h - label_h - 0.04

        placed = None
        for lane in lane_order:
            ry = clamp(by + lane*lane_step, y_lo, y_hi)
            cand = (bx, ry, bx + LABEL_W_IN, ry + label_h)
            if all(not rects_intersect(cand, r) for r in placed_rects):
                placed = cand; placed_rects.append(cand); break

        if placed is None:
            # fallback just above center
            ry = clamp(y_center - label_h - 0.02, y_lo, y_hi)
            bx = clamp(bx, x_min, x_max - LABEL_W_IN)
            placed = (bx, ry, bx + LABEL_W_IN, ry + label_h)
            placed_rects.append(placed)

        x0,y0,x1,y1 = placed
        tb = slide.shapes.add_textbox(IN(x0), IN(y0), IN(x1-x0), IN(y1-y0))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = label_text; p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER if prefer_above else PP_ALIGN.LEFT
        tf.vertical_anchor = MSO_ANCHOR.TOP

def plot_milestones(slide, df_page, groups, year, geom):
    # collect items per row, but keep month grouping for label mode choice
    per_row_by_month = {}  # (row_idx, month_idx) -> list of items
    for _, r in df_page.iterrows():
        dt = pd.to_datetime(r["Milestone Date"]).date()
        if dt.year != year: continue
        month_idx = dt.month - 1
        dim = monthrange(dt.year, dt.month)[1]
        frac = (dt.day - 1)/(dim - 1) if dim > 1 else 0.0
        x = geom['dates_left'] + (month_idx + frac) * geom['month_w']

        grp = f"{r['Type']}\n({r['Workstream']})"
        try:
            i = groups.index(grp)
        except ValueError:
            continue
        y_center = geom['first_row_top'] + i*geom['row_h'] + geom['row_h']/2.0

        # draw marker
        is_major = str(r.get("Milestone Type","")).strip().casefold() == "major"
        size = STAR_SIZE_IN if is_major else CIRCLE_SIZE_IN
        half = size/2.0
        shp = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT if is_major else MSO_SHAPE.OVAL,
                                     IN(x - half), IN(y_center - half), IN(size), IN(size))
        status = clean_text(r.get("Milestone Status",""))
        shp.fill.solid(); shp.fill.fore_color.rgb = STATUS_COLORS.get(status, RGBColor(128,128,128))
        shp.line.color.rgb = BLACK

        key = (i, month_idx)
        per_row_by_month.setdefault(key, []).append({
            'x': x,
            'y_center': y_center,
            'title': clean_text(r.get("Milestone Title",""))
        })

    # place labels per row & month bucket
    for (row_idx, month_idx), items in per_row_by_month.items():
        place_labels_for_row(slide, items, geom, month_idx)

# -----------------------
# SLIDE BUILDER
# -----------------------
def build_slide(prs, df_page, year, page_no, total_pages, title_prefix="TM‑US Roadmap"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_txt = f"{title_prefix} — {year}"
    if total_pages > 1:
        title_txt += f"  (Page {page_no}/{total_pages})"
    add_title(slide, title_txt)
    add_full_width_legend(slide)

    # groups for this page (unique ordered)
    groups = (
        df_page[["Type","Workstream"]]
        .drop_duplicates()
        .apply(lambda s: f"{s['Type']}\n({s['Workstream']})", axis=1)
        .tolist()
    )

    # build editable table (fills height), and get geometry for shapes
    geom = build_table(slide, year, groups)

    # add dotted green today line (if same year)
    add_today_line(slide, year, geom)

    # place milestones + labels
    plot_milestones(slide, df_page, groups, year, geom)

# -----------------------
# MAIN
# -----------------------
def build_ppt(input_xlsx: str, output_pptx: str, title_prefix="TM‑US Roadmap"):
    df = pd.read_excel(input_xlsx)

    # normalize
    df["Type"] = df["Type"].map(clean_text)
    df["Workstream"] = df["Workstream"].map(clean_text)
    df["Milestone Title"] = df["Milestone Title"].map(clean_text)
    df["Milestone Date"] = pd.to_datetime(df["Milestone Date"])
    df["year"] = df["Milestone Date"].dt.year

    # sort order
    df["Type_key"] = df["Type"].str.casefold()
    df["Type_bucket"] = df["Type"].map(type_bucket)
    df["Work_key"] = df["Workstream"].str.casefold()
    df = df.sort_values(by=["Type_bucket","Type_key","Work_key","Milestone Date"], kind="stable")

    # create presentation
    prs = Presentation()
    prs.slide_width  = IN(SLIDE_W_IN)
    prs.slide_height = IN(SLIDE_H_IN)

    for year in sorted(df["year"].dropna().unique()):
        df_year = df[df["year"] == year].copy()

        # ordered groups for pagination (unique Type/Workstream)
        groups_all = (
            df_year[["Type","Workstream","Type_bucket","Type_key","Work_key"]]
            .drop_duplicates()
            .sort_values(by=["Type_bucket","Type_key","Work_key"], kind="stable")
            .apply(lambda r: f"{r['Type']}\n({r['Workstream']})", axis=1)
            .tolist()
        )

        # paginate in exact chunks of 20
        pages = [groups_all[i:i+ROWS_PER_SLIDE] for i in range(0, len(groups_all), ROWS_PER_SLIDE)] or [[]]
        total_pages = len(pages)

        for page_no, group_slice in enumerate(pages, start=1):
            # filter rows for this page’s groups, keep within-page order stable
            order_index = {g:i for i,g in enumerate(group_slice)}
            df_page = df_year[df_year.apply(lambda r: f"{r['Type']}\n({r['Workstream']})" in group_slice, axis=1)].copy()
            df_page["_gidx"] = df_page.apply(lambda r: order_index[f"{r['Type']}\n({r['Workstream']})"], axis=1)
            df_page = df_page.sort_values(by=["_gidx","Milestone Date"], kind="stable").drop(columns="_gidx")

            build_slide(prs, df_page, year, page_no, total_pages, title_prefix)

    prs.save(output_pptx)
    print("Saved:", output_pptx)

# -----------------------
# Run (example)
# -----------------------
if __name__ == "__main__":
    INPUT_XLSX = r"C:\path\to\your\Roadmap_Input_Sheet.xlsx"  # change me
    OUTPUT_PPTX = r"C:\path\to\your\Roadmap.pptx"             # change me
    build_ppt(INPUT_XLSX, OUTPUT_PPTX)