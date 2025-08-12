# -*- coding: utf-8 -*-
"""
Roadmap PPT builder — full editable table + independent shapes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE

import pandas as pd
import numpy as np
import re
from datetime import datetime, date
from calendar import monthrange

# ------------------------------------------------------------
# 0) CONFIG
# ------------------------------------------------------------
IN_XLSX    = r"C:\path\to\Roadmap_Input.xlsx"          # <-- change
OUT_PPTX   = r"C:\path\to\Roadmap.pptx"                # <-- change

SLIDE_W_IN = 20.0
SLIDE_H_IN = 9.0

LEFT_PAD_IN   = 0.09
RIGHT_PAD_IN  = 0.09
TOP_PAD_IN    = 0.25
LEGEND_H_IN   = 0.45
HEADER_H_IN   = 1.0
BOTTOM_PAD_IN = 0.09

TYPE_COL_W_IN = 1.6
WORK_COL_W_IN = 2.8
MAX_ROWS_PER_SLIDE = 20

# shape sizes
CIRCLE_D_IN = 0.30
STAR_D_IN   = 0.40
LABEL_W_IN  = 2.4
LABEL_H_IN  = 0.45

# colors
BLUE_HDR   = RGBColor(91,155,213)
WHITE      = RGBColor(255,255,255)
TEXT_DARK  = RGBColor(0,0,0)
MONTH_ODD  = RGBColor(224,242,255)
MONTH_EVEN = RGBColor(190,220,240)
NAVY       = RGBColor(0,0,128)
GREEN_TOD  = RGBColor(0,176,80)

STATUS_COLORS = {
    "on track":   RGBColor(0,176,80),
    "at risk":    RGBColor(255,192,0),
    "off track":  RGBColor(255,0,0),
    "complete":   RGBColor(0,112,192),
    "tbc":        RGBColor(191,191,191),
}

# ------------------------------------------------------------
# 1) HELPERS
# ------------------------------------------------------------
def clean(s):
    if pd.isna(s): return ""
    return str(s).strip().replace("\n"," ").casefold()

def type_bucket(s):
    s = clean(s)
    m = re.match(r"([a-z0-9]+)", s)
    return m.group(1) if m else s

def build_groups_for_year(df_year_sorted):
    grp_series = df_year_sorted.apply(lambda r: f"{r['Type']}\n{r['Workstream']}", axis=1)
    groups = pd.Series(grp_series).drop_duplicates().tolist()
    return groups

def slice_pages(groups, max_rows=MAX_ROWS_PER_SLIDE):
    for i in range(0, len(groups), max_rows):
        yield i//max_rows + 1, groups[i:i+max_rows]

# ------------------------------------------------------------
# 2) TABLE BUILDER
# ------------------------------------------------------------
def build_full_table(slide, groups, year):
    rows = int(len(groups) + 1)  
    cols = int(14)              

    left_in  = LEFT_PAD_IN
    right_in = SLIDE_W_IN - RIGHT_PAD_IN
    total_w  = right_in - left_in

    type_w = TYPE_COL_W_IN
    work_w = WORK_COL_W_IN
    months_w = max(0.01, total_w - (type_w + work_w))
    month_w  = months_w / 12.0

    top_in = TOP_PAD_IN + LEGEND_H_IN + 0.10
    avail_h = SLIDE_H_IN - top_in - BOTTOM_PAD_IN
    header_h = HEADER_H_IN
    body_h   = max(0.01, avail_h - header_h)
    body_rows = max(1, rows-1)
    row_h = body_h / body_rows

    tbl_shape = slide.shapes.add_table(
        rows, cols,
        Inches(left_in),
        Inches(top_in),
        Inches(total_w),
        Inches(header_h + row_h*(rows-1))
    )
    tbl = tbl_shape.table
    tbl.style = 'Table Grid'

    tbl.columns[0].width = Inches(type_w)
    tbl.columns[1].width = Inches(work_w)
    for c in range(12):
        tbl.columns[2+c].width = Inches(month_w)

    tbl.rows[0].height = Inches(header_h)
    for r in range(1, rows):
        tbl.rows[r].height = Inches(row_h)

    hdr_cells = tbl.rows[0].cells
    hdr_cells[0].text = "Type"
    hdr_cells[1].text = "Workstream"
    for m_idx in range(12):
        dt = date(year, m_idx+1, 1)
        hdr_cells[2+m_idx].text = dt.strftime("%b %y")

    # ---- CHANGE 1: perfect center header text ----
    for c in range(cols):
        p = hdr_cells[c].text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        hdr_cells[c].text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        hdr_cells[c].margin_top = 0
        hdr_cells[c].margin_bottom = 0
        hdr_cells[c].margin_left = 0
        hdr_cells[c].margin_right = 0
        for run in p.runs: run.font.size = Pt(18)
        hdr_cells[c].fill.solid()
        hdr_cells[c].fill.fore_color.rgb = BLUE_HDR
        for run in p.runs: run.font.color.rgb = WHITE

    # body rows
    for r, grp in enumerate(groups, start=1):
        fill_rgb = MONTH_ODD if (r % 2 == 1) else MONTH_EVEN
        for c in range(cols):
            cell = tbl.rows[r].cells[c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_rgb
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            cell.margin_top = 0
            cell.margin_bottom = 0
            cell.margin_left = 0
            cell.margin_right = 0
            p.font.size = Pt(15)
            p.font.color.rgb = TEXT_DARK

        # ---- CHANGE 2: perfect center for Type & Workstream cells ----
        t, w = grp.split("\n", 1)
        for ci, text in [(0, t), (1, w)]:
            tbl.rows[r].cells[ci].text = text
            tbl.rows[r].cells[ci].text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            tbl.rows[r].cells[ci].margin_top = 0
            tbl.rows[r].cells[ci].margin_bottom = 0
            tbl.rows[r].cells[ci].margin_left = 0
            tbl.rows[r].cells[ci].margin_right = 0
            tbl.rows[r].cells[ci].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    left_months  = left_in + type_w + work_w
    right_months = left_months + months_w
    y_tops   = [ (top_in + header_h + row_h*(r-1)) for r in range(1, rows) ]
    y_centers = [ y + row_h/2.0 for y in y_tops ]  

    geom = {
        "left_in": left_in,
        "right_in": right_in,
        "top_in": top_in,
        "header_h": header_h,
        "row_h": row_h,
        "left_months_in": left_months,
        "right_months_in": right_months,
        "month_w_in": month_w,
        "y_centers_in": y_centers,
    }
    return tbl_shape, tbl, geom

# ------------------------------------------------------------
# 3) LINES & TODAY
# ------------------------------------------------------------
def draw_row_center_lines(slide, geom, row_count):
    x1 = Inches(geom["left_months_in"])
    x2 = Inches(geom["right_months_in"])
    for y_in in geom["y_centers_in"][:row_count]:
        y = Inches(y_in)
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, x1, y, x2, y
        )
        ln.line.fill.solid()
        ln.line.fill.fore_color.rgb = NAVY
        ln.line.width = Pt(0.5)

def add_today_line_if_same_year(slide, year, geom):
    today = date.today()
    if today.year != year:
        return
    days_in = monthrange(today.year, today.month)[1]
    month_idx = today.month - 1
    day_frac = (today.day - 1) / days_in
    xpos_in = geom["left_months_in"] + (month_idx + day_frac) * geom["month_w_in"]

    ln = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(xpos_in), Inches(geom["top_in"]),
        Inches(xpos_in), Inches(geom["top_in"] + geom["row_h"]*len(geom["y_centers_in"]))
    )
    ln.line.fill.solid()
    ln.line.fill.fore_color.rgb = GREEN_TOD
    ln.line.width = Pt(2)
    ln.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT

# ------------------------------------------------------------
# 4) MILESTONES & LABELS (unchanged)
# ------------------------------------------------------------
# [Keep your existing milestone plotting code here without changes]

# ------------------------------------------------------------
# 5) LEGEND (unchanged)
# ------------------------------------------------------------
# [Keep your existing legend code here without changes]

# ------------------------------------------------------------
# 6) SLIDE BUILDER & MAIN (unchanged)
# ------------------------------------------------------------
# [Keep your existing slide building + main function code here without changes]