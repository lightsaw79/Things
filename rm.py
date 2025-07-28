from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import pandas as pd
from datetime import datetime
from calendar import monthrange

# 1) Read and prep your data
data_path = r"Roadmap_Input_Sheet.xlsx"
df = pd.read_excel(data_path)
df["Milestone Date"] = pd.to_datetime(df["Milestone Date"])
df["Year"] = df["Milestone Date"].dt.year
df["Group"] = df["Type"] + "\n(" + df["Workstream"] + ")"

# 2) Create the presentation
prs = Presentation()
prs.slide_width  = Inches(20)
prs.slide_height = Inches(9)

# 3) One slide per calendar year
for year in sorted(df["Year"].unique()):
    df_year = df[df["Year"] == year].copy()
    groups = df_year["Group"].unique().tolist()

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 3A) Year title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(2), Inches(0.5))
    tf = title.text_frame
    tf.text = str(year)
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # 3B) Legend items
    legend_items = [
        ("Major Milestone", MSO_SHAPE.STAR_5_POINT, RGBColor(0,176,80)),
        ("On Track",       MSO_SHAPE.OVAL,        RGBColor(0,176,80)),
        ("At Risk",        MSO_SHAPE.OVAL,        RGBColor(255,192,0)),
        ("Off Track",      MSO_SHAPE.OVAL,        RGBColor(255,0,0)),
        ("Complete",       MSO_SHAPE.OVAL,        RGBColor(0,112,192)),
        ("TBC",            MSO_SHAPE.OVAL,        RGBColor(191,191,191)),
    ]
    sidebar_w     = Inches(4)
    type_col_w    = Inches(1.5)
    work_col_w    = sidebar_w - type_col_w
    header_h      = Inches(1)
    legend_height = Inches(0.6)
    legend_top    = Inches(0.2)
    legend_left   = Inches(0.5)
    slot_w        = (prs.slide_width - 1.0*Inches(0.5)) / len(legend_items)

    for i, (lbl, shp_type, clr) in enumerate(legend_items):
        x = legend_left + i * slot_w
        # shape
        m = slide.shapes.add_shape(
            shp_type,
            x + Inches(0.1),
            legend_top,
            Inches(0.3),
            Inches(0.3),
        )
        m.fill.solid()
        m.fill.fore_color.rgb = clr
        m.line.fill.background()
        # label
        tb = slide.shapes.add_textbox(
            x + Inches(0.5),
            legend_top,
            slot_w - Inches(0.5),
            Inches(0.3),
        )
        p = tb.text_frame.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    # 3C) Compute chart geometry
    top_margin    = legend_top + legend_height + Inches(0.2)
    bottom_margin = Inches(0.5)
    chart_w       = prs.slide_width - sidebar_w - Inches(0.5)
    chart_h       = prs.slide_height - top_margin - bottom_margin
    col_count     = 12
    row_count     = len(groups)
    col_w         = chart_w  / col_count
    row_h         = chart_h  / row_count
    left_origin   = sidebar_w + Inches(0.25)
    top_origin    = top_margin

    # 3D) Draw month header cells
    for i, m in enumerate(pd.date_range(f"{year}-01-01", f"{year}-12-01", freq="MS")):
        cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left_origin + i*col_w,
            top_origin,
            col_w,
            header_h,
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(91,155,213)
        cell.line.fill.background()
        cell.line.width = Pt(0.95)

        tf = cell.text_frame
        tf.text = m.strftime("%b %y")
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = RGBColor(255,255,255)
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 3E) Draw sidebar header “Type” / “Workstream”
    th1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.25),
        top_origin,
        type_col_w,
        header_h,
    )
    th1.fill.solid()
    th1.fill.fore_color.rgb = RGBColor(91,155,213)
    th1.line.fill.background()
    th1.text_frame.text = "Type"
    p = th1.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER
    th1.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    th2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.25)+type_col_w,
        top_origin,
        work_col_w,
        header_h,
    )
    th2.fill.solid()
    th2.fill.fore_color.rgb = RGBColor(91,155,213)
    th2.line.fill.background()
    th2.text_frame.text = "Workstream"
    p = th2.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER
    th2.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 3F) Draw the alternating sidebar rows
    for r, grp in enumerate(groups):
        y = top_origin + header_h + r*row_h
        bg = RGBColor(242,242,242) if r%2==0 else RGBColor(190,220,240)
        # Type cell
        c1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.25),
            y,
            type_col_w,
            row_h,
        )
        c1.fill.solid(); c1.fill.fore_color.rgb = bg
        c1.line.fill.background(); c1.line.width = Pt(0.5)
        tf = c1.text_frame; tf.text = grp.split("\n")[0]
        p = tf.paragraphs[0]; p.font.size = Pt(10); p.alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Workstream cell
        c2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.25)+type_col_w,
            y,
            work_col_w,
            row_h,
        )
        c2.fill.solid(); c2.fill.fore_color.rgb = bg
        c2.line.fill.background(); c2.line.width = Pt(0.5)
        tf = c2.text_frame; tf.text = grp.split("\n")[1].strip("()")
        p = tf.paragraphs[0]; p.font.size = Pt(10); p.alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 3G) Draw grid cells & horizontal center‐lines
    for i in range(col_count):
        for r in range(row_count):
            x = left_origin + i*col_w
            y = top_origin + header_h + r*row_h
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y, col_w, row_h
            )
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255,255,255)
            cell.line.fill.background(); cell.line.width = Pt(0.5)
    # horizontal lines:
    for r in range(row_count+1):
        y = top_origin + header_h + r*row_h
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            left_origin, y,
            left_origin+chart_w, y
        )
        ln.line.fill.solid()
        ln.line.fill.fore_color.rgb = RGBColor(128,128,128)
        ln.line.width = Pt(0.5)

    # 3H) Plot each milestone
    status_colors = {
        "On Track":   RGBColor(0,176,80),
        "At Risk":    RGBColor(255,192,0),
        "Off Track":  RGBColor(255,0,0),
        "Complete":   RGBColor(0,112,192),
        "TBC":        RGBColor(191,191,191),
    }
    shape_map = {"Regular": MSO_SHAPE.OVAL, "Major": MSO_SHAPE.STAR_5_POINT}

    for _, row in df_year.iterrows():
        xi = row["Milestone Date"].month - 1
        days_in_month = monthrange(year, row["Milestone Date"].month)[1]
        day_frac = (row["Milestone Date"].day - 1) / (days_in_month - 1)
        x = left_origin + (xi + day_frac)*col_w + col_w/2 - Inches(0.15)
        yi = groups.index(row["Group"])
        y = top_origin + header_h + yi*row_h + row_h/2 - Inches(0.15)

        shp = slide.shapes.add_shape(
            shape_map[row["Milestone Type"]],
            x, y,
            Inches(0.3), Inches(0.3)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = status_colors[row["Milestone Status"]]
        shp.line.fill.background()

        lbl = slide.shapes.add_textbox(
            x + Inches(0.4),
            y - Inches(0.1),
            Inches(2), Inches(0.3)
        )
        tf = lbl.text_frame
        tf.text = row["Milestone Title"]
        tf.paragraphs[0].font.size = Pt(10)

    # 3I) Draw “today” dotted vertical line
    today = datetime.today().date()
    if today.year == year:
        days_in_month = monthrange(year, today.month)[1]
        day_frac = (today.day - 1) / (days_in_month - 1)
        xpos = left_origin + (today.month-1 + day_frac)*col_w
        today_ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            xpos, top_origin + header_h,
            xpos, top_origin + header_h + chart_h
        )
        today_ln.line.fill.solid()
        today_ln.line.fill.fore_color.rgb = RGBColor(0,176,80)
        today_ln.line.width = Pt(1)
        today_ln.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT

# 4) Save the PPTX
output_path = r"Roadmap_by_Year.pptx"
prs.save(output_path)
print(f"Saved → {output_path}")